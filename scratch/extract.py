import sys
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_pdf(path):
    print(f"--- Extracting PDF: {path} ---")
    try:
        reader = PdfReader(path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            print(f"Page {i+1}:\n{text}\n")
    except Exception as e:
        print(f"Error reading PDF: {e}")

def extract_pptx(path):
    print(f"--- Extracting PPTX: {path} ---")
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
            print()
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    extract_pdf(r"d:\abi\skynet\output\Skynet Professional Pitch Deck.pdf")
    extract_pptx(r"d:\abi\skynet\output\skynet-nidhi-prayas-deck-polished.pptx")
