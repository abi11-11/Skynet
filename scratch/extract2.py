import sys
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_pdf(path, outfile):
    outfile.write(f"--- Extracting PDF: {path} ---\n")
    try:
        reader = PdfReader(path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            outfile.write(f"Page {i+1}:\n{text}\n\n")
    except Exception as e:
        outfile.write(f"Error reading PDF: {e}\n")

def extract_pptx(path, outfile):
    outfile.write(f"--- Extracting PPTX: {path} ---\n")
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            outfile.write(f"Slide {i+1}:\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    outfile.write(shape.text + "\n")
            outfile.write("\n")
    except Exception as e:
        outfile.write(f"Error reading PPTX: {e}\n")

if __name__ == "__main__":
    with open(r"d:\abi\skynet\scratch\extract_output_utf8.txt", "w", encoding="utf-8") as f:
        extract_pdf(r"d:\abi\skynet\output\Skynet Professional Pitch Deck.pdf", f)
        extract_pptx(r"d:\abi\skynet\output\skynet-nidhi-prayas-deck-polished.pptx", f)
