import sys
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# Use blank layout (index 6) to build custom layouts
blank_slide_layout = prs.slide_layouts[6]

def add_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 22)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_bg(slide)
    
    # Accent bar
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(0.1)
    height = Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 200, 255)
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(8), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(160, 160, 160)
    p2.font.name = 'Arial'

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_bg(slide)
    
    # Accent line top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 200, 255)
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    
    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.font.name = 'Arial'
        p.space_after = Pt(20)
        p.level = 0

# --- Presentation Content ---

add_title_slide(
    "Skynet", 
    "Verifiable Drone Spray Coverage for Tamil Nadu Farmers.\nHardware-verified, farmer-trusted."
)

add_content_slide(
    "The Problem",
    [
        "Farmers pay for drone spraying but cannot verify if the entire field was covered.",
        "Lack of proof leads to disputes and low trust between farmers and drone pilots.",
        "Manual labour shortages make drones essential, but current solutions waste chemicals and increase costs.",
        "Smallholder farmers are excluded due to the lack of affordable and trustworthy Drone-as-a-Service (DaaS) platforms."
    ]
)

add_content_slide(
    "Our Solution: PCVM Gen-1",
    [
        "Hardware-verified proof of drone spray coverage.",
        "Centimetre-accurate GNSS track logging directly on the drone.",
        "Tamper-evident, offline-first logs ensuring data integrity in rural areas.",
        "Geo-verified coverage maps mapped against actual field boundaries.",
        "Exportable 'proof packets' providing verifiable coverage percentages to the farmer."
    ]
)

add_content_slide(
    "Innovation & Differentiation",
    [
        "Low-cost hardware combined with advanced analytics for verifiable proof.",
        "Designed specifically for rural operations with an offline-first, tamper-evident architecture.",
        "Automated boundary-alignment and gap-detection systems.",
        "Acts as a foundational trust layer for DaaS pilots and Farmer Producer Organizations (FPOs)."
    ]
)

add_content_slide(
    "Market Opportunity",
    [
        "Target Market: Coimbatore → Western Tamil Nadu → Delta Belt.",
        "DaaS Pricing Band: ₹300–700 per acre.",
        "Highly scalable entry clusters: 500 farmers × 4 sprays.",
        "Single pilot cluster GMV is highly profitable and sets up immediate regional expansion."
    ]
)

add_content_slide(
    "Business Model",
    [
        "Revenue Stream 1: Per-acre service fee for verified spraying.",
        "Revenue Stream 2: Dedicated contracts with FPOs and large estates.",
        "Revenue Stream 3: PCVM hardware rental/sale post-commercialization.",
        "Revenue Stream 4: Aggregated precision data services (Phase 2)."
    ]
)

add_content_slide(
    "Competitive Landscape",
    [
        "Competitors rely on manual pilot reports or expensive enterprise software.",
        "Skynet Differentiator 1: Hardware-verified cryptographic proof.",
        "Skynet Differentiator 2: FPO-first go-to-market strategy.",
        "Skynet Differentiator 3: Built for rural India with deep offline UX capabilities."
    ]
)

add_content_slide(
    "Roadmap",
    [
        "Phase 1: Build and field-validate the PCVM Gen-1 prototype.",
        "Phase 2: Pilot batch (20-30 units) deployment and beta bookings.",
        "Phase 3: Secure FPO contracts and validate unit economics.",
        "Phase 4: Scale operations and launch advanced data services."
    ]
)

add_content_slide(
    "Traction & Validation",
    [
        "Comprehensive market research and Product Requirements Document (PRD) finalized.",
        "Hardware prototypes actively entering development phase.",
        "Planned field validation includes over 10 field trials and 20+ direct feedback loops.",
        "Establishing MOUs with regional pilot clusters."
    ]
)

add_content_slide(
    "The Team",
    [
        "A multi-disciplinary team combining hardware, firmware, and business strategy.",
        "Expertise in Agriculture technology and Drone Operations.",
        "Deep understanding of the Tamil Nadu farming landscape.",
        "Committed to building scalable, trust-driven platforms for rural economies."
    ]
)

add_title_slide(
    "Thank You",
    "Skynet - Making every acre accountable.\nContact the Lead Innovator for investment or partnership opportunities."
)

prs.save(r'd:\abi\skynet\output\skynet-pitch-deck-professional.pptx')
print("Presentation generated successfully!")
