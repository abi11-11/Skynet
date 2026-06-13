from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pathlib import Path

ASSET_DIR = Path('output/planning-artifacts/grants/nidhi-prayas/assets')
OUTPUT = Path('output/skynet-nidhi-prayas-deck-professional.pptx')

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
navy = RGBColor(10, 50, 80)
green = RGBColor(0, 145, 110)
light_green = RGBColor(215, 245, 230)
text_dark = RGBColor(22, 30, 39)
white = RGBColor(255, 255, 255)

def add_footer(slide, text):
    tx = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.5), Inches(12.5), Inches(0.4))
    tf = tx.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = white

    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(8.0), Inches(2.0))
    tf = title.text_frame
    tf.text = 'Skynet'
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = navy

    subtitle = slide.shapes.add_textbox(Inches(0.75), Inches(2.3), Inches(9.0), Inches(1.3))
    tf2 = subtitle.text_frame
    tf2.text = 'Verifiable Drone Spray Coverage for Tamil Nadu Farmers'
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(40, 70, 115)

    details = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(9.0), Inches(1.5))
    tf3 = details.text_frame
    tf3.text = 'NIDHI-PRAYAS Prototype Grant Application'
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = text_dark
    for text in [
        'Innovator / Startup: Skynet',
        'Location: Coimbatore, Tamil Nadu',
        'Grant sought: ₹10,00,000',
        'Duration: 6–8 months'
    ]:
        p = tf3.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = text_dark

    if (ASSET_DIR / 'logo.png').exists():
        slide.shapes.add_picture(str(ASSET_DIR / 'logo.png'), Inches(9.5), Inches(0.9), width=Inches(3.0))
    add_footer(slide, 'Skynet • NIDHI-PRAYAS grant pitch deck • May 2026')


def add_content_slide(title, bullets, notes=None, image=None, image_pos=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(7.5), Inches(5.8))
    tf = body.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = navy
    tf.paragraphs[0].space_after = Pt(14)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = text_dark
        p.space_after = Pt(6)

    if notes:
        note_box = slide.shapes.add_textbox(Inches(8.4), Inches(5.0), Inches(4.2), Inches(1.3))
        ntf = note_box.text_frame
        ntf.text = notes
        ntf.paragraphs[0].font.size = Pt(12)
        ntf.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)

    if image and (ASSET_DIR / image).exists():
        x, y, width = image_pos if image_pos else (Inches(8.5), Inches(1.0), Inches(4.3))
        slide.shapes.add_picture(str(ASSET_DIR / image), x, y, width=width)

    add_footer(slide, 'Skynet • NIDHI-PRAYAS • May 2026')


def add_problem_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(1.0)).text_frame.text = 'The Problem'
    title = slide.shapes[0].text_frame.paragraphs[0]
    title.font.size = Pt(36)
    title.font.bold = True
    title.font.color.rgb = navy

    issues = [
        ('No proof of coverage', 'Farmers cannot verify spray execution, generating disputes and limiting trust.'),
        ('Manual labour shortage', 'TN has acute farm labour gaps; drone spraying is essential.'),
        ('Waste & cost risk', 'Over/under-spray reduces yield and raises input costs.'),
        ('Smallholder access', 'Farmers need affordable DaaS, not expensive drone ownership.')
    ]
    left = Inches(0.6)
    top = Inches(1.6)
    for i, (hdr, desc) in enumerate(issues):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + (i%2)*Inches(6.4), top + (i//2)*Inches(2.1), Inches(6.0), Inches(1.9))
        box.fill.solid()
        box.fill.fore_color.rgb = light_green
        box.line.color.rgb = green
        tx = box.text_frame
        tx.text = hdr
        p = tx.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = text_dark
        p = tx.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(45, 55, 70)

    add_footer(slide, 'Skynet • NIDHI-PRAYAS • Problem statement')


def add_two_column_slide(title, left_items, right_image=None, right_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.9))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = navy

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.6), Inches(5.4))
    ltf = left_box.text_frame
    ltf.text = left_items[0]
    ltf.paragraphs[0].font.size = Pt(24)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = text_dark
    for item in left_items[1:]:
        p = ltf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(45, 55, 70)
        p.space_after = Pt(8)

    if right_image and (ASSET_DIR / right_image).exists():
        slide.shapes.add_picture(str(ASSET_DIR / right_image), Inches(7.5), Inches(1.5), width=Inches(5.0))
    elif right_note:
        note_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5.0), Inches(3.5))
        ntf = note_box.text_frame
        ntf.text = right_note
        ntf.paragraphs[0].font.size = Pt(16)
        ntf.paragraphs[0].font.color.rgb = RGBColor(45, 55, 70)

    add_footer(slide, 'Skynet • NIDHI-PRAYAS • May 2026')


def add_table_slide(title, headers, rows, image=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.9))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(36)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = navy

    rows_count = len(rows) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows_count, cols, Inches(0.6), Inches(1.6), Inches(7.5), Inches(4.5)).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = king_green if False else light_green
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(13)

    if image and (ASSET_DIR / image).exists():
        slide.shapes.add_picture(str(ASSET_DIR / image), Inches(8.6), Inches(1.6), width=Inches(4.0))

    add_footer(slide, 'Skynet • NIDHI-PRAYAS • May 2026')


# Build slides
add_title_slide()
add_problem_slide()
add_two_column_slide(
    'Our Solution',
    [
        'Proof-of-Coverage Verification Module (PCVM) Gen-1',
        '• Captures GPS flight track on spray runs',
        '• Stores tamper-evident offline logs',
        '• Aligns coverage to farmer KMZ boundaries',
        '• Produces exportable GPX/GeoJSON proof packets'
    ],
    right_image='coverage_diagram.png'
)
add_content_slide(
    'Innovation & Differentiation',
    [
        'Low-cost hardware + analytics for verifiable service proof',
        'Offline-first field design with tamper-evident logging',
        'Boundary alignment and gap detection for precision',
        'Trust layer for DaaS, FPOs, and smallholder pilots'
    ]
)
add_content_slide(
    'Prototype Focus',
    [
        'Field-ready PCVM hardware logger for agricultural spray drones',
        'GNSS tracker, MCU edge compute, rugged enclosure, field mount',
        'Offline proof packet generation; easy drone integration',
        'Built for 6–8 month NIDHI-PRAYAS prototype validation'
    ]
)
add_content_slide(
    'Market Opportunity',
    [
        'Tamil Nadu entry: Coimbatore → Western TN → Delta paddy belt',
        'DaaS pricing band: ₹300–700/acre',
        'Year-1 cluster economics: 500 farmers × 4 sprays × ₹400 = ₹8L GMV',
        'Growth path: FPO contracts, service bundles, PCVM licensing'
    ]
)
add_content_slide(
    'Business Model',
    [
        'Per-acre service fee for verified spray delivery',
        'Seasonal bulk pricing for FPOs and estates',
        'PCVM rental/sale to pilot partners post-commercialization',
        'Data services upsell (NDVI, monitoring) in Phase 2'
    ]
)
add_content_slide(
    'Competitive Edge',
    [
        'Hardware-verified proof, not just a booking app',
        'Designed for rural offline conditions and smallholders',
        'Built around FPO/CHC GTM and local pilot trust',
        'Aligns directly with NIDHI-PRAYAS prototype criteria'
    ]
)
add_content_slide(
    'Roadmap',
    [
        'Month 1: Design freeze, BOM lock, procurement kickoff',
        'Months 2–3: Lab prototype build and GPS log validation',
        'Months 4–5: Drone integration and tethered tests',
        'Months 6–7: 10+ field trials, farmer feedback, field hardening',
        'Month 8: Final report, commercialization package, DST deliverables'
    ]
)
add_two_column_slide(
    'Budget Overview',
    [
        '₹10,00,000 budget aligned to NIDHI-PRAYAS heads',
        'Sensors, electronics, and mechanical integration',
        'Field trials, PRAYAS Shala prototyping, and compliance',
        'Contingency for iteration and supply-chain risk'
    ],
    right_image='budget_chart.png'
)
add_content_slide(
    'Team & Execution',
    [
        'Lead Innovator: product + grant ownership',
        'Hardware/IoT: sensor integration and prototype build',
        'Agriculture advisor: TN farming and FPO relationships',
        'Drone operations: licensed pilot support and test execution',
        'Business/GTM: local market launch and FPO engagement'
    ]
)
add_content_slide(
    'Traction & Validation',
    [
        'Market research and PRD completed',
        'Pre-grant status: prototype not yet built',
        'Grant objective: build and validate PCVM Gen-1',
        'Field validation: 10+ demo plots and 20+ stakeholder sessions'
    ]
)
add_content_slide(
    'IP & Compliance',
    [
        'IP to vest with applicant per PRAYAS/DST terms',
        'Indian citizenship and NOC requirement captured',
        'DGCA-aligned drone operations with licensed pilots',
        'Prototype tests in controlled field conditions only'
    ]
)
add_content_slide(
    'Commercialization Path',
    [
        'Months 9–12: PCVM pilot batch (20–30 units) and beta partners',
        'Months 12–18: FPO contracts and early unit economics',
        'Year 3+: regional expansion, data services, PCVM licensing',
        'Objective: TRL 6 field validation converted to market-ready offering'
    ]
)
add_content_slide(
    'The Ask',
    [
        'Request ₹10,00,000 for PCVM Gen-1 prototype and validation',
        'Deliverables: field-tested physical prototype, proof reports, DST submission pack',
        'Outcome: verified DaaS trust infrastructure for TN agriculture',
        'Aligned with NIDHI-PRAYAS priority sectors: Agriculture + IoT'
    ]
)
add_content_slide(
    'Appendix & Submission Notes',
    [
        'Completed application form and supporting ID documents',
        'Team & Eligibility details in team-and-eligibility.md',
        'Prototype spec and BOM in prototype-spec.md',
        'Detailed line-item budget in nidhi-prayas-budget.csv'
    ]
)

prs.save(OUTPUT)
print('Created', OUTPUT)
