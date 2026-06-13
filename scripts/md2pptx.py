"""Simple Markdown to PPTX converter for the NIDHI-PRAYAS deck.

Usage:
    python md2pptx.py ../output/planning-artifacts/grants/nidhi-prayas/skynet-nidhi-prayas-pitch-deck-final.md output/skynet-nidhi-prayas-deck.pptx

Requires: python-pptx
"""
import sys
from pptx import Presentation
from pptx.util import Pt

TEMPLATE_TITLE_SLIDE_LAYOUT = 0
TEMPLATE_CONTENT_SLIDE_LAYOUT = 1

def md_to_slides(md_text):
    # Split slides on lines that are only '---'
    parts = []
    current = []
    for line in md_text.splitlines():
        if line.strip() == '---':
            if current:
                parts.append('\n'.join(current).strip())
                current = []
        else:
            current.append(line)
    if current:
        parts.append('\n'.join(current).strip())
    return parts

def slide_title_and_body(slide_md):
    lines = [l for l in slide_md.splitlines() if l.strip()!='']
    if not lines:
        return ('','')
    # Title: first line that starts with '#'
    title = ''
    body_lines = []
    if lines[0].startswith('#'):
        title = lines[0].lstrip('#').strip()
        body_lines = lines[1:]
    else:
        # fallback: first non-empty line as title
        title = lines[0].strip()
        body_lines = lines[1:]
    body = '\n'.join(body_lines)
    return (title, body)

def create_pptx(md_path, pptx_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        md_text = f.read()
    slides_md = md_to_slides(md_text)
    prs = Presentation()
    for i, s in enumerate(slides_md):
        title, body = slide_title_and_body(s)
        if i == 0:
            slide_layout = prs.slide_layouts[TEMPLATE_TITLE_SLIDE_LAYOUT]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            if slide.placeholders:
                # put body into subtitle placeholder if available
                try:
                    slide.placeholders[1].text = body
                except Exception:
                    pass
        else:
            slide_layout = prs.slide_layouts[TEMPLATE_CONTENT_SLIDE_LAYOUT]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            # Add body textbox
            left = Pt(40)
            top = Pt(120)
            width = Pt(860)
            height = Pt(420)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            for paragraph in body.split('\n'):
                p = tf.add_paragraph()
                p.text = paragraph
                p.font.size = Pt(18)
    prs.save(pptx_path)

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print('Usage: python md2pptx.py <input.md> <output.pptx>')
        sys.exit(2)
    md_path = sys.argv[1]
    pptx_path = sys.argv[2]
    create_pptx(md_path, pptx_path)
    print('Wrote', pptx_path)
