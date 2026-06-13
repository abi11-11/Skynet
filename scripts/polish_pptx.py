from pptx import Presentation
from pptx.util import Inches, Pt
import os

INPUT = 'output/skynet-nidhi-prayas-deck.pptx'
OUTPUT = 'output/skynet-nidhi-prayas-deck-polished.pptx'
ASSET_DIR = 'output/planning-artifacts/grants/nidhi-prayas/assets'

prs = Presentation(INPUT)

# add logo to title slide (slide 0)
logo_path = os.path.join(ASSET_DIR, 'logo.png')
if os.path.exists(logo_path):
    slide = prs.slides[0]
    slide.shapes.add_picture(logo_path, prs.slide_width - Inches(3.5), Inches(0.2), width=Inches(3))

# helper: find slide by title keyword
def find_slide_by_keyword(keyword):
    for s in prs.slides:
        for shape in s.shapes:
            if hasattr(shape, 'text') and keyword.lower() in shape.text.lower():
                return s
    return None

# add coverage diagram to slide with 'Proof-of-Coverage' or 'Our Solution'
target = find_slide_by_keyword('Proof-of-Coverage') or find_slide_by_keyword('Our Solution')
if target:
    cov = os.path.join(ASSET_DIR, 'coverage_diagram.png')
    if os.path.exists(cov):
        target.shapes.add_picture(cov, Inches(0.5), Inches(2.0), width=Inches(9))

# add budget chart to slide with 'Budget' keyword
budget_slide = find_slide_by_keyword('Budget') or find_slide_by_keyword('Grant Utilization')
if budget_slide:
    chart = os.path.join(ASSET_DIR, 'budget_chart.png')
    if os.path.exists(chart):
        budget_slide.shapes.add_picture(chart, Inches(6.5), Inches(2.0), width=Inches(3))

# increase body font sizes for readability
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    run.font.size = Pt(18)
                except Exception:
                    pass

prs.save(OUTPUT)
print('Wrote', OUTPUT)
